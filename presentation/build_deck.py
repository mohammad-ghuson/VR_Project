#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Build the Arabic committee presentation for the Swinging Paint Bucket (VR) project.

One slide spec (px on a 1280×720 grid, authored right-to-left) is emitted twice:
  presentation/deck_ar.pptx   — the deliverable (native RTL, Noto Sans / Noto Sans Arabic)
  presentation/preview.html   — a pixel-twin for headless-browser visual QA
                                (orange numbered badges = click/animation order)

Conventions follow the team's approved decks (AlgorithmsAdult, football-simulation):
light theme, flat rounded cards, kicker + title header, centered title slide,
Arabic paragraphs rtl+right, Latin/technical islands LTR, one idea per click,
slide numbers bottom-left.

    python3 presentation/build_deck.py
"""

import html as html_mod
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

HERE = Path(__file__).parent

# ── theme: navy (report continuity) + paint-orange accent, light surfaces ─────
NAVY = "#1F3B63"
NAVY_D = "#152A47"
NAVY_T = "#EDF2F9"
NAVY_BR = "#C6D4E8"
ON_NAVY = "#C9D8ED"      # light text on navy backgrounds
ON_NAVY2 = "#9FB6D6"
ON_NAVY3 = "#DDE7F5"
ORANGE = "#C2410C"
ORANGE_T = "#FFF3EA"
ORANGE_BR = "#F0C9AC"
INK = "#14181B"
INK2 = "#4A5258"
INK3 = "#6B7480"
BORDER = "#E2E5E8"
SURFACE = "#F6F8F7"
WHITE = "#FFFFFF"
# paint colours (diagrams only)
PAINT_B = "#2E86AB"
PAINT_B2 = "#4FA3C7"
PAINT_R = "#E4572E"
PAINT_Y = "#F6C445"
PAINT_P = "#7B4B94"

SANS = "Noto Sans"
AR = "Noto Sans Arabic"

PX = 9525                # EMU per px (96 dpi on a 1280×720 grid)
N_SLIDES = 12

# ── spec model ────────────────────────────────────────────────────────────────
DECK = []


def slide(bg=WHITE):
    s = {"bg": bg, "shapes": []}
    DECK.append(s)
    return s


def rect(s, x, y, w, h, fill=None, line=None, lw=1.5, radius=10, g=None):
    s["shapes"].append(dict(kind="rect", x=x, y=y, w=w, h=h,
                            fill=fill, line=line, lw=lw, radius=radius, g=g))


def oval(s, x, y, w, h, fill=None, line=None, lw=1.5, g=None):
    s["shapes"].append(dict(kind="oval", x=x, y=y, w=w, h=h,
                            fill=fill, line=line, lw=lw, g=g))


def seg(s, x1, y1, x2, y2, color, lw=2.5, arrow=False, g=None):
    s["shapes"].append(dict(kind="seg", x1=x1, y1=y1, x2=x2, y2=y2,
                            color=color, lw=lw, arrow=arrow, g=g))


def P(text, size=14, bold=False, color=INK, align="r", rtl=True, spacing=1.18):
    """One paragraph with a single run (the common case)."""
    return dict(runs=[(text, dict(size=size, bold=bold, color=color))],
                align=align, rtl=rtl, spacing=spacing)


def text(s, x, y, w, h, paras, anchor="t", g=None):
    s["shapes"].append(dict(kind="text", x=x, y=y, w=w, h=h,
                            paras=paras, anchor=anchor, g=g))


# ── composite helpers ─────────────────────────────────────────────────────────
def header(s, kicker, title, accent=NAVY):
    text(s, 60, 38, 1160, 26, [P(kicker, size=13, bold=True, color=accent)])
    text(s, 60, 66, 1160, 50, [P(title, size=27, bold=True, color=INK)])


def slideno(s, n, color=INK3, align="l"):
    text(s, 60, 688, 200, 20,
         [P(f"{n} / {N_SLIDES}", size=9, color=color, align=align, rtl=False)])


def card(s, x, y, w, h, title, body, fill=SURFACE, line=BORDER,
         tcolor=NAVY, tsize=15.5, bsize=12.5, bcolor=INK2, trtl=True, g=None):
    rect(s, x, y, w, h, fill=fill, line=line, g=g)
    text(s, x + 10, y, w - 20, h,
         [P(title, size=tsize, bold=True, color=tcolor, align="ctr", rtl=trtl),
          P(body, size=bsize, color=bcolor, align="ctr")],
         anchor="m", g=g)


# ══ S1 — title ════════════════════════════════════════════════════════════════
s = slide(bg=NAVY)
text(s, 140, 92, 1000, 30, [P("جامعة دمشق — كلية الهندسة المعلوماتية",
                              size=15, color=ON_NAVY, align="ctr")])
text(s, 140, 126, 1000, 26, [P("مادة الحقائق الافتراضية · 2025 / 2026",
                               size=12, color=ON_NAVY2, align="ctr")])
text(s, 90, 208, 1100, 100, [P("دلو الطلاء المتأرجح",
                               size=42, bold=True, color=WHITE, align="ctr", spacing=1.2)])
text(s, 190, 300, 900, 26, [P("Swinging Paint Bucket", size=14, color=ON_NAVY2,
                              align="ctr", rtl=False)])
text(s, 140, 344, 1000, 34, [P("دلوٌ يتأرجح، وطلاءٌ يتدفق، ولوحةٌ تُرسَم أمامكم — دون أي محرك فيزيائي جاهز",
                               size=16, color=ON_NAVY3, align="ctr")])
rect(s, 565, 404, 150, 3, fill=ON_NAVY2, radius=0)
text(s, 90, 432, 1100, 28,
     [P("محمد مروان العيشات  ·  محمد فرحان غصن  ·  محمد سمير الجاعوني",
        size=14, bold=True, color=WHITE, align="ctr")])
text(s, 90, 464, 1100, 28,
     [P("محمد ياسر الدخل الله  ·  محمد قاسم السحوم  ·  علاء محمد الشريف",
        size=14, bold=True, color=WHITE, align="ctr")])
text(s, 190, 512, 900, 26, [P("بإشراف: المهندس خالد إسماعيل",
                              size=13, color=ON_NAVY, align="ctr")])
text(s, 540, 683, 200, 20, [P(f"1 / {N_SLIDES}", size=9, color=ON_NAVY2,
                              align="ctr", rtl=False)])

# ══ S2 — الفكرة في جملة واحدة ════════════════════════════════════════════════
s = slide()
header(s, "ما المشروع؟", "الفكرة في جملة واحدة")
text(s, 650, 152, 570, 120,
     [P("بنينا مختبرًا فيزيائيًا افتراضيًا: دلو طلاءٍ يتأرجح كالنواس، ويرسم لوحة فنية حقيقية.",
        size=21, bold=True, spacing=1.3)])
for i, b in enumerate([
    "كل قانون حركة مكتوب بأيدينا — لا RigidBody ولا Colliders",
    "كل مدخل فيزيائي بيد المستخدم أثناء التشغيل",
    "كل تجربة تُحفَظ: صورة + تقرير + مقارنة",
]):
    text(s, 650, 292 + i * 42, 570, 34, [P("•  " + b, size=16)], g=i + 1)
rect(s, 850, 442, 370, 46, fill=NAVY_T, line=NAVY_BR, g=4)
text(s, 850, 442, 370, 46, [P("Unity · C# · SPH — from scratch", size=14, bold=True,
                              color=NAVY, align="ctr", rtl=False)], anchor="m", g=4)

# drawn scene: rope + bucket + falling drops + canvas with a paint trail
rect(s, 60, 160, 550, 400, fill=WHITE, line=BORDER)
seg(s, 200, 205, 470, 205, INK3, lw=3)                       # ceiling
oval(s, 328, 198, 14, 14, fill=INK2)                          # anchor
seg(s, 335, 205, 428, 332, INK2, lw=2.5)                      # rope (swung)
rect(s, 396, 330, 64, 52, fill=SURFACE, line=INK2, lw=2, radius=6)   # bucket
rect(s, 402, 340, 52, 16, fill=PAINT_B, radius=4)             # paint inside
oval(s, 418, 396, 10, 10, fill=PAINT_B)                       # drops arc
oval(s, 394, 430, 10, 10, fill=PAINT_B)
oval(s, 362, 462, 11, 11, fill=PAINT_B)
oval(s, 326, 490, 12, 12, fill=PAINT_B)
rect(s, 120, 514, 430, 24, fill=WHITE, line=INK2, lw=2, radius=4)    # canvas
oval(s, 196, 520, 10, 10, fill=PAINT_P)                       # trail on canvas
oval(s, 246, 519, 12, 12, fill=PAINT_R)
oval(s, 298, 518, 14, 13, fill=PAINT_B)
oval(s, 352, 519, 12, 12, fill=PAINT_Y)
text(s, 60, 570, 550, 22, [P("المشهد كاملًا: تأرجحٌ في الأعلى — لوحةٌ تتلوّن في الأسفل",
                             size=11, color=INK3, align="ctr")])
slideno(s, 2)

# ══ S3 — سلسلة الأحداث ═══════════════════════════════════════════════════════
s = slide()
header(s, "كيف يعمل؟", "من التأرجح إلى اللوحة — ست حلقات")
steps = [
    ("يتأرجح الدلو", "معادلة النواس + تخامد الهواء"),
    ("يتموّج الطلاء", "مائع SPH حقيقي داخل الدلو"),
    ("يخرج من الفتحة", "قطر الفتحة بيد المستخدم"),
    ("تسقط القطرات", "جاذبية ومقاومة هواء"),
    ("ترتطم باللوحة", "تصادم تحليلي مع المستوي"),
    ("تتراكم لوحةً", "بصمة لونية فوق بصمة"),
]
for i, (t, d) in enumerate(steps):
    x = 60 + (5 - i) * 195           # rightmost card first (RTL flow)
    oval(s, x + 72, 176, 36, 36, fill=NAVY, g=i + 1)
    text(s, x + 72, 176, 36, 36, [P(str(i + 1), size=15, bold=True, color=WHITE,
                                    align="ctr", rtl=False)], anchor="m", g=i + 1)
    card(s, x, 228, 180, 150, t, d, tsize=14.5, bsize=11.5, g=i + 1)
text(s, 60, 452, 1160, 32, [P("كل حلقة في السلسلة فيزياء كتبناها بأنفسنا — لا شيء جاهز",
                              size=16, bold=True, color=INK, align="ctr")], g=7)
text(s, 60, 492, 1160, 28, [P("وستشاهدون السلسلة كاملة تعمل في العرض الحي",
                              size=13, color=INK3, align="ctr")], g=7)
slideno(s, 3)

# ══ S4 — نجمة المشروع: SPH ═══════════════════════════════════════════════════
s = slide()
header(s, "قلب المشروع", "الطلاء مائعٌ حقيقي — لا خدعة بصرية", accent=ORANGE)
text(s, 680, 150, 540, 92,
     [P("كل قطرةٍ جسيمٌ في مائع SPH — الطريقة المستخدمة في مؤثرات الأفلام — بنيناها من الصفر.",
        size=16, spacing=1.35)])
chips = [
    "الجسيم يشعر بجيرانه: الضغط يمنع التكدّس",
    "اللزوجة احتكاكٌ داخلي: طلاءٌ سميك أو سائلٌ رشيق",
    "الدلو يتحرك من تحته ← فيتموّج وينسكب طبيعيًا",
]
for i, c in enumerate(chips):
    rect(s, 680, 252 + i * 56, 540, 46, fill=ORANGE_T, line=ORANGE_BR, g=i + 1)
    text(s, 692, 252 + i * 56, 516, 46, [P(c, size=14, bold=True, color=INK,
                                           align="ctr")], anchor="m", g=i + 1)
text(s, 680, 432, 540, 30, [P("الاسم العلمي: SPH — الهيدروديناميكا الجسيمية الملسّاء (Müller 2003)",
                              size=12.5, color=INK2)], g=4)

# drawn bucket-of-particles
rect(s, 60, 165, 560, 330, fill=WHITE, line=BORDER)
seg(s, 230, 230, 230, 400, INK2, lw=3)                        # left wall
seg(s, 450, 230, 450, 400, INK2, lw=3)                        # right wall
seg(s, 230, 400, 318, 400, INK2, lw=3)                        # bottom (hole gap)
seg(s, 362, 400, 450, 400, INK2, lw=3)
pts = [(248, 372), (274, 380), (302, 370), (330, 378), (358, 372), (386, 380),
       (414, 371), (260, 344), (292, 340), (324, 348), (356, 342), (390, 346),
       (276, 314), (312, 310), (348, 316), (382, 312), (300, 286), (340, 284)]
for j, (px_, py_) in enumerate(pts):
    oval(s, px_, py_, 17, 17, fill=PAINT_B if j % 3 else PAINT_B2)
oval(s, 332, 418, 12, 12, fill=PAINT_B)                       # drops out of hole
oval(s, 336, 446, 12, 12, fill=PAINT_B)
oval(s, 330, 472, 13, 13, fill=PAINT_B)
text(s, 372, 414, 90, 20, [P("الفتحة", size=11, color=INK3)])
text(s, 60, 502, 560, 22, [P("جسيمات تتدافع وتنسكب من الفتحة — لا مقاطع جاهزة",
                             size=11, color=INK3, align="ctr")])
slideno(s, 4)

# ══ S5 — الأداء: برهان O(n) ══════════════════════════════════════════════════
s = slide()
header(s, "الأداء", "أسرع بأربع مرات ونصف — والبرهان على الشاشة")
stats = [
    (847, "≈ 14 ms", "الشبكة الذكية O(n) — عند 3000 جسيم", 1),
    (453, "≈ 65 ms", "الطريقة الساذجة O(n²) — العمل نفسه", 2),
    (60, "PASS", "فحص التطابق: الفيزياء لم تتغيّر", 3),
]
for x, big, label, g_ in stats:
    rect(s, x, 150, 373, 112, fill=SURFACE, line=BORDER, g=g_)
    text(s, x + 12, 150, 349, 112,
         [P(big, size=30, bold=True, color=NAVY if g_ != 2 else ORANGE,
            align="ctr", rtl=False),
          P(label, size=12.5, color=INK2, align="ctr")], anchor="m", g=g_)
text(s, 660, 300, 560, 24, [P("الشبكة المكانية O(n)", size=12.5, bold=True, color=NAVY)], g=4)
rect(s, 1050, 328, 170, 22, fill=NAVY, radius=6, g=4)
text(s, 660, 362, 560, 24, [P("القوة الغاشمة O(n²)", size=12.5, bold=True, color=ORANGE)], g=4)
rect(s, 450, 390, 770, 22, fill=ORANGE, radius=6, g=4)
text(s, 60, 330, 360, 60, [P("الزمن الحقيقي المقاس للبحث عن الجيران",
                             size=12, color=INK3, spacing=1.3)], g=4)
text(s, 60, 452, 1160, 32,
     [P("زرٌّ واحد في الواجهة يبدّل بين الطريقتين — والميلي ثانية تُقاس أمامكم مباشرة",
        size=16, bold=True, color=NAVY, align="ctr")], g=5)
text(s, 60, 492, 1160, 28,
     [P("لماذا لا نطلب أسرع من O(n)؟ لأن كل جسيم يجب أن يتحدّث مرة كل خطوة — الخطّي هو الحد الأمثل",
        size=13, color=INK3, align="ctr")], g=5)
slideno(s, 5)

# ══ S6 — كيف بُنيت: المكوّنات ════════════════════════════════════════════════
s = slide()
header(s, "تحت الغطاء", "أربعة مكوّنات — لكلٍّ مهمة واحدة")
comps = [
    ("Bucket", "يحسب التأرجح من معادلة النواس — رياضيات صرفة"),
    ("SphFluid", "يحاكي آلاف جسيمات الطلاء: ضغط ولزوجة وتموّج"),
    ("PaintCanvas", "يستقبل القطرات ويرسم اللوحة ويحفظها"),
    ("UIControlPanel", "يربط كل منزلقٍ بالفيزياء — حيًّا"),
]
for i, (t, d) in enumerate(comps):
    x = 72 + (3 - i) * 288
    card(s, x, 216, 272, 148, t, d, tsize=16, bsize=12.5, trtl=False, g=i + 1)
text(s, 60, 428, 1160, 30, [P("فصلٌ تام بين المسؤوليات: كل مكوّن يُفحَص ويُشرَح وحده",
                              size=15, color=INK, align="ctr")], g=5)
text(s, 60, 468, 1160, 32,
     [P("القيد الأساسي محترمٌ بالكامل: لا RigidBody · لا Collider · لا مكتبات جاهزة",
        size=15, bold=True, color=INK, align="ctr")], g=6)
text(s, 60, 508, 1160, 28, [P("حتى الحدود مصادمة تحليلية: أسطوانة للدلو ومستوٍ للّوحة — بالمعادلات",
                              size=13, color=INK3, align="ctr")], g=6)
slideno(s, 6)

# ══ S7 — كل شيء بيد المستخدم ═════════════════════════════════════════════════
s = slide()
header(s, "التحكم", "غيّر الفيزياء… تتغيّر اللوحة")
ctrls = [
    (938, 195, "الدلو والطلاء", "الكمية · قطر الفتحة · اللون · اللزوجة"),
    (660, 195, "الحركة", "زاوية البداية · السرعة · عدد التأرجحات · طول الحبل"),
    (938, 311, "البيئة", "الجاذبية · مقاومة الهواء · الرطوبة · الاحتكاك"),
    (660, 311, "اللوحة", "الأبعاد · الميلان حتى 60° · قماش / خشب / معدن / ورق"),
]
for gi, (x, y, t, d) in enumerate(ctrls):
    card(s, x, y, 262, 100, t, d, tcolor=INK, tsize=15, bsize=11.5, g=gi + 1)
text(s, 660, 445, 540, 28, [P("كل مدخلٍ في ملف المشروع… منزلقٌ حيٌّ أمامكم",
                              size=13, color=INK2)], g=5)

# drawn control panel: sliders + colour swatches
rect(s, 60, 170, 550, 330, fill=WHITE, line=BORDER)
sliders = [("الجاذبية", 330), ("اللزوجة", 210), ("الرطوبة", 390), ("طول الحبل", 262)]
for r_i, (lbl, knob) in enumerate(sliders):
    y = 232 + r_i * 54
    text(s, 452, y - 11, 130, 22, [P(lbl, size=11.5, color=INK2)])
    seg(s, 112, y, 436, y, BORDER, lw=4)
    oval(s, knob, y - 8, 16, 16, fill=NAVY)
text(s, 452, 437, 130, 22, [P("اللون", size=11.5, color=INK2)])
for c_i, c in enumerate((PAINT_R, PAINT_B, PAINT_Y, PAINT_P)):
    oval(s, 160 + c_i * 46, 430, 24, 24, fill=c)
text(s, 60, 508, 550, 22, [P("لوحة التحكم الحقيقية: أكثر من عشرين منزلقًا وزرًّا حيًّا",
                             size=11, color=INK3, align="ctr")])
slideno(s, 7)

# ══ S8 — المخرجات ════════════════════════════════════════════════════════════
s = slide()
header(s, "المخرجات", "كل تجربة توثّق نفسها")
outs = [
    (847, "صورة PNG", "اللوحة النهائية تُحفَظ إلى ملف بنقرة واحدة"),
    (453, "تقرير نصي", "كل المدخلات + زمن الحركة + عدد الآثار + مساحة الانتشار"),
    (60, "مقارنة حية", "التجربة السابقة مقابل الحالية — جنبًا إلى جنب"),
]
for gi, (x, t, d) in enumerate(outs):
    card(s, x, 165, 373, 150, t, d, tsize=17, bsize=12.5, g=gi + 1)
text(s, 60, 380, 1160, 32,
     [P("وعلى الشاشة دائمًا: قيم التجربة الجارية لحظةً بلحظة",
        size=15, color=INK, align="ctr")], g=4)
text(s, 60, 430, 1160, 34,
     [P("المخرجات السبعة المطلوبة في ملف المشروع — منفّذة كلها",
        size=17, bold=True, color=NAVY, align="ctr")], g=5)
text(s, 60, 472, 1160, 28,
     [P("عرض ثلاثي الأبعاد · رسم حي · لوحة نهائية · حفظ · قيم · مقارنة · تقرير",
        size=13, color=INK3, align="ctr")], g=5)
slideno(s, 8)

# ══ S9 — العرض الحي (interstitial) ═══════════════════════════════════════════
s = slide(bg=NAVY_D)
text(s, 90, 188, 1100, 30, [P("العرض الحي", size=14, bold=True, color=ON_NAVY2, align="ctr")])
text(s, 90, 248, 1100, 90, [P("لنشاهد الدلو يرسم", size=44, bold=True, color=WHITE, align="ctr")])
text(s, 90, 368, 1100, 36,
     [P("لوحة بعدة ألوان  ·  رطوبة أعلى فبقعٌ أوسع  ·  برهان O(n) أمامكم",
        size=16, color=ON_NAVY3, align="ctr")])
rect(s, 490, 560, 300, 46, fill=None, line=ON_NAVY2, lw=1.5)
text(s, 490, 560, 300, 46, [P("Unity · SampleScene", size=14, bold=True,
                              color=WHITE, align="ctr", rtl=False)], anchor="m")
slideno(s, 9, color=ON_NAVY2)

# ══ S10 — الصدق العلمي ═══════════════════════════════════════════════════════
s = slide()
header(s, "خيارات نمذجة صادقة", "بسيطٌ وصحيح خيرٌ من معقّدٍ خاطئ")
rect(s, 660, 170, 560, 120, fill=SURFACE, line=BORDER, g=1)
text(s, 672, 180, 536, 102,
     [P("وزن الدلو", size=14.5, bold=True, color=NAVY),
      P("دور النواس البسيط مستقلٌّ عن الكتلة — فلا منزلق وزنٍ زائف. فيزياء، لا نسيان.",
        size=13, color=INK2, spacing=1.3)], g=1)
rect(s, 660, 306, 560, 120, fill=SURFACE, line=BORDER, g=2)
text(s, 672, 316, 536, 102,
     [P("الحبل", size=14.5, bold=True, color=NAVY),
      P("قضيبٌ صلب اليوم — والحبل المرن موثَّق كعمل مستقبلي، لا مُدَّعى.",
        size=13, color=INK2, spacing=1.3)], g=2)
rect(s, 60, 170, 560, 256, fill=NAVY_T, line=NAVY_BR, g=3)
text(s, 84, 190, 512, 220,
     [P("«كل تبسيطٍ مدوَّنٌ كتابةً", size=20, bold=True, color=NAVY, align="ctr", spacing=1.4),
      P("ومبرَّرٌ فيزيائيًا.»", size=20, bold=True, color=NAVY, align="ctr", spacing=1.4)],
     anchor="m", g=3)
text(s, 60, 470, 1160, 28,
     [P("القادم: عدة دلاء ترسم معًا · لعبة تعليمية · تصدير فيديو · لوحات بالذكاء الاصطناعي",
        size=13, color=INK3, align="ctr")], g=4)
slideno(s, 10)

# ══ S11 — الخلاصة ════════════════════════════════════════════════════════════
s = slide()
text(s, 90, 205, 1100, 26, [P("الخلاصة", size=13, bold=True, color=NAVY, align="ctr")])
text(s, 150, 255, 980, 170,
     [P("«لا شيء في المشهد مزيّف: من تأرجح الحبل إلى آخر قطرةٍ على اللوحة، "
        "كلُّ سلوكٍ ينبع من معادلةٍ كتبناها بأيدينا.»",
        size=25, bold=True, color=NAVY, align="ctr", spacing=1.4)], g=1)
text(s, 190, 445, 900, 28, [P("التفاصيل والقيود موثّقة بصدقٍ في التقرير النهائي.",
                              size=13, color=INK2, align="ctr")], g=2)
slideno(s, 11)

# ══ S12 — الفريق والشكر ══════════════════════════════════════════════════════
s = slide()
text(s, 90, 118, 1100, 56, [P("شكرًا لكم", size=30, bold=True, color=INK, align="ctr")])
names = ["محمد مروان العيشات", "محمد فرحان غصن", "محمد سمير الجاعوني",
         "محمد ياسر الدخل الله", "محمد قاسم السحوم", "علاء محمد الشريف"]
for i, name in enumerate(names):
    row, col = divmod(i, 3)
    x = 88 + (2 - col) * 376
    y = 226 + row * 84
    rect(s, x, y, 352, 64, fill=SURFACE, line=BORDER)
    text(s, x, y, 352, 64, [P(name, size=14.5, bold=True, align="ctr")], anchor="m")
text(s, 190, 412, 900, 30, [P("بإشراف: المهندس خالد إسماعيل", size=14, color=INK2, align="ctr")])
text(s, 90, 480, 1100, 70, [P("أسئلتكم؟", size=34, bold=True, color=NAVY, align="ctr")])
slideno(s, 12)


# ══ pptx emitter ══════════════════════════════════════════════════════════════
def rgb(hexs):
    return RGBColor(int(hexs[1:3], 16), int(hexs[3:5], 16), int(hexs[5:7], 16))


def _style_run(r, st):
    f = r.font
    f.name = SANS
    f.size = Pt(st["size"])
    f.bold = st["bold"]
    f.color.rgb = rgb(st["color"])
    rPr = r._r.get_or_add_rPr()
    for tag in ("a:latin", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", SANS if tag == "a:latin" else AR)


ALIGN = {"r": PP_ALIGN.RIGHT, "l": PP_ALIGN.LEFT, "ctr": PP_ALIGN.CENTER}


def _animate(slide_obj, groups):
    """Click-appear animation via timing-XML injection (the ResQRules recipe).

    groups: click-ordered list; each entry = list of pptx shapes that appear
    together on that click. Shapes not listed stay visible from the start.
    """
    from lxml import etree
    groups = [g for g in groups if g]
    if not groups:
        return
    nid = [1]

    def i_():
        nid[0] += 1
        return nid[0]

    ns = ('xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
          'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"')
    clicks = []
    for g in groups:
        effects = []
        for j, sp in enumerate(g):
            node_t = "clickEffect" if j == 0 else "withEffect"
            effects.append(f"""
              <p:par><p:cTn id="{i_()}" presetID="1" presetClass="entr" presetSubtype="0" fill="hold" grpId="0" nodeType="{node_t}">
                <p:stCondLst><p:cond delay="0"/></p:stCondLst>
                <p:childTnLst><p:set><p:cBhvr>
                  <p:cTn id="{i_()}" dur="1" fill="hold"><p:stCondLst><p:cond delay="0"/></p:stCondLst></p:cTn>
                  <p:tgtEl><p:spTgt spid="{sp.shape_id}"/></p:tgtEl>
                  <p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>
                </p:cBhvr><p:to><p:strVal val="visible"/></p:to></p:set></p:childTnLst>
              </p:cTn></p:par>""")
        clicks.append(f"""
          <p:par><p:cTn id="{i_()}" fill="hold">
            <p:stCondLst><p:cond delay="indefinite"/></p:stCondLst>
            <p:childTnLst><p:par><p:cTn id="{i_()}" fill="hold">
              <p:stCondLst><p:cond delay="0"/></p:stCondLst>
              <p:childTnLst>{''.join(effects)}</p:childTnLst>
            </p:cTn></p:par></p:childTnLst>
          </p:cTn></p:par>""")
    xml = f"""<p:timing {ns}><p:tnLst><p:par>
      <p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot"><p:childTnLst>
        <p:seq concurrent="1" nextAc="seek">
          <p:cTn id="{i_()}" dur="indefinite" nodeType="mainSeq"><p:childTnLst>
            {''.join(clicks)}
          </p:childTnLst></p:cTn>
          <p:prevCondLst><p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:prevCondLst>
          <p:nextCondLst><p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:nextCondLst>
        </p:seq>
      </p:childTnLst></p:cTn></p:par></p:tnLst></p:timing>"""
    slide_obj._element.append(etree.fromstring(xml.encode()))


def emit_pptx(out_path):
    prs = Presentation()
    prs.slide_width = Emu(1280 * PX)
    prs.slide_height = Emu(720 * PX)
    blank = prs.slide_layouts[6]

    for spec in DECK:
        sl = prs.slides.add_slide(blank)
        sl.background.fill.solid()
        sl.background.fill.fore_color.rgb = rgb(spec["bg"])
        pending = {}
        for sh in spec["shapes"]:
            created = None
            k = sh["kind"]
            if k in ("rect", "oval"):
                mso = MSO_SHAPE.OVAL if k == "oval" else (
                    MSO_SHAPE.ROUNDED_RECTANGLE if sh.get("radius", 10) else MSO_SHAPE.RECTANGLE)
                box = sl.shapes.add_shape(mso, Emu(sh["x"] * PX), Emu(sh["y"] * PX),
                                          Emu(sh["w"] * PX), Emu(sh["h"] * PX))
                if k == "rect" and sh.get("radius", 10):
                    box.adjustments[0] = min(0.5, sh["radius"] / min(sh["w"], sh["h"]))
                if sh["fill"]:
                    box.fill.solid()
                    box.fill.fore_color.rgb = rgb(sh["fill"])
                else:
                    box.fill.background()
                if sh["line"]:
                    box.line.color.rgb = rgb(sh["line"])
                    box.line.width = Pt(sh["lw"])
                else:
                    box.line.fill.background()
                box.shadow.inherit = False
                created = box
            elif k == "seg":
                from pptx.enum.shapes import MSO_CONNECTOR
                c = sl.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                            Emu(sh["x1"] * PX), Emu(sh["y1"] * PX),
                                            Emu(sh["x2"] * PX), Emu(sh["y2"] * PX))
                c.line.color.rgb = rgb(sh["color"])
                c.line.width = Pt(sh["lw"])
                c.shadow.inherit = False
                if sh["arrow"]:
                    ln = c.line._get_or_add_ln()
                    tail = ln.makeelement(qn("a:tailEnd"),
                                          {"type": "triangle", "w": "med", "len": "med"})
                    ln.append(tail)
                created = c
            elif k == "text":
                tb = sl.shapes.add_textbox(Emu(sh["x"] * PX), Emu(sh["y"] * PX),
                                           Emu(sh["w"] * PX), Emu(sh["h"] * PX))
                tf = tb.text_frame
                tf.word_wrap = True
                tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE if sh["anchor"] == "m" else MSO_ANCHOR.TOP
                for pi, para in enumerate(sh["paras"]):
                    p = tf.paragraphs[0] if pi == 0 else tf.add_paragraph()
                    p.alignment = ALIGN[para["align"]]
                    p.line_spacing = para["spacing"]
                    if para["rtl"]:
                        p._p.get_or_add_pPr().set("rtl", "1")
                    for txt, st in para["runs"]:
                        _style_run(p.add_run(), st)
                        p.runs[-1].text = txt
                created = tb
            if sh.get("g") and created is not None:
                pending.setdefault(sh["g"], []).append(created)
        _animate(sl, [pending[k] for k in sorted(pending)])
    prs.save(out_path)
    print(f"wrote {out_path}")


# ══ html emitter (pixel twin for QA) ══════════════════════════════════════════
def emit_html(out_path):
    css_align = {"r": "right", "l": "left", "ctr": "center"}
    parts = ["""<!doctype html><html><head><meta charset="utf-8">
<link href="https://fonts.googleapis.com/css2?family=Noto+Sans:wght@400;600;700&family=Noto+Sans+Arabic:wght@400;600;700&display=swap" rel="stylesheet">
<style>
 body{background:#3a3f44;margin:0;font-family:'Noto Sans','Noto Sans Arabic',sans-serif;}
 .slide{position:relative;width:1280px;height:720px;margin:24px auto;overflow:hidden;}
 .slide *{box-sizing:border-box;}
 p{margin:0;}
</style></head><body>"""]
    for spec in DECK:
        parts.append(f'<div class="slide" style="background:{spec["bg"]}">')
        svg_lines = []
        badges = []
        for sh in spec["shapes"]:
            k = sh["kind"]
            if sh.get("g"):
                bx = sh.get("x", sh.get("x1", 0))
                by = sh.get("y", sh.get("y1", 0))
                badges.append(
                    f'<div style="position:absolute;left:{bx - 9}px;top:{by - 9}px;'
                    f'width:20px;height:20px;border-radius:50%;background:#C2410C;'
                    f'color:#fff;font-size:11px;font-weight:700;display:flex;'
                    f'align-items:center;justify-content:center;opacity:.9;z-index:9">'
                    f'{sh["g"]}</div>')
            if k == "rect":
                st = (f'position:absolute;left:{sh["x"]}px;top:{sh["y"]}px;'
                      f'width:{sh["w"]}px;height:{sh["h"]}px;'
                      f'border-radius:{sh.get("radius", 10)}px;')
                if sh["fill"]:
                    st += f'background:{sh["fill"]};'
                if sh["line"]:
                    st += f'border:{sh["lw"]}px solid {sh["line"]};'
                parts.append(f'<div style="{st}"></div>')
            elif k == "oval":
                st = (f'position:absolute;left:{sh["x"]}px;top:{sh["y"]}px;'
                      f'width:{sh["w"]}px;height:{sh["h"]}px;border-radius:50%;')
                if sh["fill"]:
                    st += f'background:{sh["fill"]};'
                if sh["line"]:
                    st += f'border:{sh["lw"]}px solid {sh["line"]};'
                parts.append(f'<div style="{st}"></div>')
            elif k == "seg":
                import math
                x1, y1, x2, y2 = sh["x1"], sh["y1"], sh["x2"], sh["y2"]
                svg_lines.append(
                    f'<line x1="{x1}" y1="{y1}" x2="{x2}" y2="{y2}" '
                    f'stroke="{sh["color"]}" stroke-width="{sh["lw"]}"/>')
                if sh["arrow"]:
                    ang = math.atan2(y2 - y1, x2 - x1)
                    L = 5 * sh["lw"]
                    pts_ = []
                    for da in (0.46, -0.46):
                        pts_.append((x2 - L * math.cos(ang + da), y2 - L * math.sin(ang + da)))
                    svg_lines.append(
                        f'<polygon points="{x2},{y2} {pts_[0][0]:.1f},{pts_[0][1]:.1f} '
                        f'{pts_[1][0]:.1f},{pts_[1][1]:.1f}" fill="{sh["color"]}"/>')
            elif k == "text":
                anchor = ("display:flex;flex-direction:column;justify-content:center;"
                          if sh["anchor"] == "m" else "")
                parts.append(f'<div style="position:absolute;left:{sh["x"]}px;top:{sh["y"]}px;'
                             f'width:{sh["w"]}px;height:{sh["h"]}px;{anchor}">')
                for para in sh["paras"]:
                    direction = "rtl" if para["rtl"] else "ltr"
                    runs = "".join(
                        f'<span style="font-size:{st["size"] * 4 / 3:.1f}px;'
                        f'font-weight:{700 if st["bold"] else 400};color:{st["color"]}">'
                        f'{html_mod.escape(txt)}</span>'
                        for txt, st in para["runs"])
                    parts.append(f'<p dir="{direction}" style="text-align:'
                                 f'{css_align[para["align"]]};line-height:{para["spacing"]}">'
                                 f'{runs}</p>')
                parts.append("</div>")
        if svg_lines:
            parts.append(
                '<svg style="position:absolute;inset:0;pointer-events:none" width="1280" height="720">'
                + "".join(svg_lines) + "</svg>")
        parts.extend(badges)
        parts.append("</div>")
    parts.append("</body></html>")
    Path(out_path).write_text("\n".join(parts), encoding="utf-8")
    print(f"wrote {out_path}")


if __name__ == "__main__":
    emit_pptx(HERE / "deck_ar.pptx")
    emit_html(HERE / "preview.html")
